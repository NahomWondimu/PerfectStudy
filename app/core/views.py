from django.shortcuts import render
#do pip install google-genai
from google import genai
from google.genai import types
import typing
import json
from rest_framework.decorators import api_view
from rest_framework.response import Response
import os
from dotenv import load_dotenv
from django.http import HttpResponse, JsonResponse
from PyPDF2 import PdfReader
import io
from pptx import Presentation
# Create your views here.

load_dotenv()

class FlashCard(typing.TypedDict):
    question: str
    answer: str


class FeynmanInput(typing.TypedDict):
    score: int
    response: str


API_KEY = os.getenv("GEMINI_API_KEY")
client = genai.Client(api_key=API_KEY)

@api_view(['GET', 'POST'])
def index(request):

    if request.method == "POST":

        text = ""
        json_data = None

        file = request.FILES.get('file')
        if str(file).endswith(".pdf"):

            text = ""
            file = request.FILES.get('file')
            reader = PdfReader(file)
            

            for i in range(len(reader.pages)):
                text += reader.pages[i].extract_text()
            

            json_data = getGeminiResponse(text)
    

        elif str(file).endswith(".pptx"):
            text = ""
            presentation = Presentation(file)
            for slide_num, slide in enumerate(presentation.slides):
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text

            json_data = getGeminiResponse(text)
    
        return JsonResponse(json_data, safe=False)
    
    return render(request, "index.html")


@api_view(['GET', 'POST'])
def feynman(request):
    
    if request.method == "POST":

        str_builder = ""

        topic = request.POST.get("topic")
        elementaryResponse = request.POST.get("elementaryResponse")
        middleResponse = request.POST.get("middleResponse")
        highResponse = request.POST.get("highResponse")

        str_builder += "Topic: " + topic
        str_builder += " Elementary School Response: " + elementaryResponse 
        str_builder += " Middle School Response: " + middleResponse
        str_builder += " High School Response: " + highResponse

        prompt = "You will give 3 scores from 0-100 on how good the responses are based on how well it can be explained to the school level chosen by the user. Additionally, you will be giving your own feedback on how they can improve their responses for each of the school levels. It will be based on this data: " + str_builder

        response = client.models.generate_content(
            model = "gemini-2.0-flash",
            contents=prompt,
            config = {
                "response_mime_type": "application/json",
                "response_schema": list[FeynmanInput]
            }
        )

        text_data = response.text
        json_data = json.loads(text_data)

        return JsonResponse(json_data, safe=False)

    else:
        return render(request, "index.html")


def getGeminiResponse(content):
    response = client.models.generate_content(

        model="gemini-2.0-flash", 
        contents="Create 10 short flashcards using content from this document: " + content,
        config = {
            'response_mime_type': 'application/json',
            'response_schema': list[FlashCard]
        }
    )

    text_data = response.text
    json_data = json.loads(text_data)

    return json_data







